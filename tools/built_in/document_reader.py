"""
DocumentReader — reads and understands PDF, Word, Excel, PowerPoint, and images.
No need to open the app — the agent reads document content directly.
"""

from __future__ import annotations

from pathlib import Path
from tools.base_tool import BaseTool
from utils.logger import get_logger

log = get_logger(__name__)


class DocumentReaderTool(BaseTool):
    name = "read_document"
    description = (
        "Read and extract text from documents: PDF, Word (.docx), Excel (.xlsx), "
        "PowerPoint (.pptx), or images with text. Returns the content as plain text."
    )
    parameters = {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "File path to read"},
            "max_chars": {"type": "integer", "default": 5000},
            "sheet": {"type": "string", "description": "Excel sheet name (optional)"},
        },
        "required": ["path"],
    }

    def run(self, path: str, max_chars: int = 5000, sheet: str = "") -> str:
        p = Path(path)
        if not p.exists():
            return f"File not found: {path}"

        ext = p.suffix.lower()
        try:
            if ext == ".pdf":
                return self._read_pdf(p, max_chars)
            elif ext in (".docx", ".doc"):
                return self._read_word(p, max_chars)
            elif ext in (".xlsx", ".xls", ".csv"):
                return self._read_excel(p, max_chars, sheet)
            elif ext in (".pptx", ".ppt"):
                return self._read_pptx(p, max_chars)
            elif ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"):
                return self._read_image(p, max_chars)
            elif ext in (".txt", ".md", ".log", ".json", ".yaml", ".yml", ".toml", ".ini", ".cfg"):
                return p.read_text(encoding="utf-8", errors="replace")[:max_chars]
            else:
                return f"Unsupported format: {ext}"
        except Exception as e:
            return f"Error reading {path}: {e}"

    def _read_pdf(self, p: Path, max_chars: int) -> str:
        try:
            import pdfplumber
            text = []
            with pdfplumber.open(p) as pdf:
                for page in pdf.pages:
                    t = page.extract_text()
                    if t:
                        text.append(t)
            return "\n".join(text)[:max_chars]
        except ImportError:
            # Fallback: pymupdf
            try:
                import fitz
                doc = fitz.open(str(p))
                return "\n".join(page.get_text() for page in doc)[:max_chars]
            except ImportError:
                return "PDF reading requires: pip install pdfplumber  OR  pip install pymupdf"

    def _read_word(self, p: Path, max_chars: int) -> str:
        try:
            from docx import Document
            doc = Document(str(p))
            text = "\n".join(para.text for para in doc.paragraphs if para.text.strip())
            # Also get table content
            for table in doc.tables:
                for row in table.rows:
                    text += "\n" + " | ".join(cell.text for cell in row.cells)
            return text[:max_chars]
        except ImportError:
            return "Word reading requires: pip install python-docx"

    def _read_excel(self, p: Path, max_chars: int, sheet: str = "") -> str:
        try:
            import openpyxl
            if p.suffix.lower() == ".csv":
                import csv
                rows = []
                with open(p, encoding="utf-8", errors="replace") as f:
                    reader = csv.reader(f)
                    for row in reader:
                        rows.append(" | ".join(row))
                return "\n".join(rows)[:max_chars]

            wb = openpyxl.load_workbook(str(p), read_only=True, data_only=True)
            ws = wb[sheet] if sheet and sheet in wb.sheetnames else wb.active
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    rows.append(" | ".join(cells))
            return f"Sheet: {ws.title}\n" + "\n".join(rows)[:max_chars]
        except ImportError:
            return "Excel reading requires: pip install openpyxl"

    def _read_pptx(self, p: Path, max_chars: int) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(str(p))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
                slides.append(f"Slide {i}:\n" + "\n".join(texts))
            return "\n\n".join(slides)[:max_chars]
        except ImportError:
            return "PowerPoint reading requires: pip install python-pptx"

    def _read_image(self, p: Path, max_chars: int) -> str:
        try:
            import pytesseract
            from PIL import Image
            img = Image.open(p)
            text = pytesseract.image_to_string(img)
            return text[:max_chars] if text.strip() else "[No text found in image]"
        except ImportError:
            return "Image OCR requires: pytesseract + Tesseract binary"
        except Exception as e:
            return f"Image read error: {e}"
